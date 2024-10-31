import os
from pptx import Presentation
import comtypes.client as client  # for Windows PowerPoint automation (to convert .ppt to .pptx)
import re

def convert_ppt_to_pptx(ppt_path):
    powerpoint = client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    ppt = powerpoint.Presentations.Open(ppt_path)
    new_path = re.sub(r"\.ppt$", ".pptx", ppt_path, flags=re.IGNORECASE)
    ppt.SaveAs(new_path, 24)  # 24 is the format type for .pptx
    ppt.Close()
    powerpoint.Quit()
    return new_path

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
    return text

def process_ppt_files(folder_path):
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        if filename.endswith(".ppt"):
            try:
                print(f"Converting {filename} to .pptx...")
                file_path = convert_ppt_to_pptx(file_path)
            except Exception as e:
                print(f"Failed to convert {filename}: {e}")
                continue
        if filename.endswith(".pptx"):
            print(f"Extracting text from {filename}...")
            text = extract_text_from_pptx(file_path)
            output_file = os.path.join(folder_path, f"{os.path.splitext(filename)[0]}.txt")
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(text)
            print(f"Text extracted to {output_file}")

# Define your folder path
folder_path = "path_to_your_folder"
process_ppt_files(folder_path)
